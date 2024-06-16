from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import matplotlib.pyplot as plt
# Data provided
data = {
    "RL": [117, 117.1, 117.2, 117.3, 117.4, 117.5, 117.6, 117.7, 117.8, 117.9, 118, 118.1],
    "Area": [1823.9244, 2128.85226, 2433.78012, 2738.70798, 3043.63584, 3348.5637, 3653.49156, 3958.41942, 4263.34728, 4568.27514, 4873.203, 5077.40871],
    "Volume": [91.19622, 197.638833, 425.770452, 684.394857, 973.512048, 1293.122025, 1643.224788, 2023.820337, 2434.908672, 2876.489793, 3348.5637, 3846.094286]
}
# Create DataFrame
df = pd.DataFrame(data)
# Create a Presentation object
prs = Presentation()
# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Sump Volume Calculation Before Monsoon"
subtitle.text = "Prepared by [Your Team's Names]\nDate: [Current Date]"
# Introduction Slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = "Overview of the project\nObjectives of the sump volume calculation"
# Methodology Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Methodology"
content.text = "Tools and techniques used: Excel, AutoCAD, contour lines\nRole of the surveyor: Tarun Panda"
# Data Presentation Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Data Presentation"
# Add Table to Data Presentation Slide
rows, cols = df.shape[0] + 1, df.shape[1] + 1
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(9.0), Inches(3.5)).table
# Set column names
table.cell(0, 1).text = "RL"
table.cell(0, 2).text = "Area of Each Contour"
table.cell(0, 3).text = "Volume"
# Add data to table
for i in range(len(df)):
    table.cell(i + 1, 1).text = str(df['RL'][i])
    table.cell(i + 1, 2).text = str(df['Area'][i])
    table.cell(i + 1, 3).text = str(df['Volume'][i])
# Create a plot
plt.figure(figsize=(10, 6))
plt.plot(df['RL'], df['Volume'], marker='o', linestyle='-', color='b')
plt.title('Volume vs RL')
plt.xlabel('RL')
plt.ylabel('Volume')
plt.grid(True)
plt.savefig('volume_vs_rl.png')
plt.close()
# Add the plot to the slide
slide.shapes.add_picture('volume_vs_rl.png', Inches(1), Inches(4.5), Inches(8), Inches(4))
# Calculation Results Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Calculation Results"
content.text = "Detailed results of the volume calculation for East and West sump\nInterpretation of the results"
# Conclusion Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = "Summary of findings\nImportance of the results for monsoon preparation"
# Acknowledgements Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Acknowledgements"
content.text = "Thanks to team members and surveyor Tarun Panda"
# Save presentation
pptx_file = "Sump_Volume_Calculation_Before_Monsoon.pptx"
prs.save(pptx_file)
pptx_file
