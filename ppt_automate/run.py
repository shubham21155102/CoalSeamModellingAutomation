
# import Presentation class 
# from pptx library 
from pptx import Presentation  
# Creating presentation object 
root = Presentation() 
# Creating slide layout 
first_slide_layout = root.slide_layouts[0]  
  
""" Ref for slide types:  
0 ->  title and subtitle 
1 ->  title and content 
2 ->  section header 
3 ->  two content 
4 ->  Comparison 
5 ->  Title only  
6 ->  Blank 
7 ->  Content with caption 
8 ->  Pic with caption 
"""
# Creating slide object to add  
# in ppt i.e. Attaching slides  
# with Presentation i.e. ppt 
slide = root.slides.add_slide(first_slide_layout) 
  
# Adding title and subtitle in  
# slide i.e. first page of slide  
slide.shapes.title.text = " Created By python-pptx"
  
# We have different formats of  
# subtitles in ppts, for simple 
# subtitle this method should  
# implemented, you can change 
# 0 to 1 for different design 
slide.placeholders[1].text = " This is 2nd way"
  
# Saving file 
root.save("Output.pptx") 
  
print("done") 